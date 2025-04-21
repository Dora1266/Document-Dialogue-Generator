import os
import io
import re
import fitz
import docx
import csv
import json
import uuid
import time
import logging
import tempfile
import argparse
import requests
import concurrent.futures
import nltk
from nltk.tokenize import sent_tokenize
from typing import List, Dict, Any, Tuple
from flask import Flask, request, jsonify, render_template, send_file
from flask_cors import CORS
from werkzeug.utils import secure_filename
from datetime import datetime

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger('document_dialogue_app')

try:
    import spacy
    import easyocr
    import numpy as np
    from PIL import Image
    from pptx import Presentation
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl import load_workbook

    try:
        nlp = spacy.load('zh_core_web_sm')
    except OSError:
        try:
            nlp = spacy.load('en_core_web_sm')
        except OSError:
            nlp = None

    try:
        nltk.data.find('tokenizers/punkt')
    except LookupError:
        nltk.download('punkt', quiet=True)

    OPTIONAL_DEPS_AVAILABLE = True
except ImportError:
    OPTIONAL_DEPS_AVAILABLE = False
    nlp = None
    logger.warning("Some optional dependencies are missing. Limited functionality available.")


def ensure_directory(directory_path):
    if not os.path.exists(directory_path):
        try:
            os.makedirs(directory_path, mode=0o755)
            logger.info(f"Created directory: {directory_path}")
        except Exception as e:
            logger.error(f"Failed to create directory {directory_path}: {e}")
            return False

    try:
        test_file = os.path.join(directory_path, ".write_test")
        with open(test_file, 'w') as f:
            f.write("test")
        os.remove(test_file)
        return True
    except Exception as e:
        logger.error(f"Directory {directory_path} is not writable: {e}")
        return False


class SemanticDocumentParser:
    def __init__(self, language_list=None, segment_size=1000, overlap=True, overlap_limit=200,
                 clean_for_ai=True, replace_whitespace=False,
                 remove_urls_emails=False, disable_ocr=False):
        if language_list is None:
            language_list = ['ch_sim', 'en']

        self.languages = language_list
        self.segment_size = segment_size
        self.overlap = overlap
        self.overlap_limit = overlap_limit
        self.nlp = nlp
        self.clean_for_ai = clean_for_ai
        self.replace_whitespace = replace_whitespace
        self.remove_urls_emails = remove_urls_emails
        self.disable_ocr = disable_ocr

        self.primary_language = 'en'
        if 'ch_sim' in language_list or 'ch_tra' in language_list:
            self.primary_language = 'zh'

        if OPTIONAL_DEPS_AVAILABLE and not self.disable_ocr:
            try:
                self.reader = easyocr.Reader(language_list)
            except Exception as e:
                logger.error(f"Failed to initialize EasyOCR: {e}")
                self.reader = None
        else:
            self.reader = None

    def parse(self, file_path, segment_size=None):
        if segment_size is not None:
            self.segment_size = segment_size

        try:
            _, ext = os.path.splitext(file_path)
            ext = ext.lower()
        except Exception as e:
            file_name = os.path.basename(file_path)
            ext = '.' + file_name.split('.')[-1].lower() if '.' in file_name else ''
            logger.warning(f"Using fallback to get file extension: {ext}, error: {str(e)}")

        if not ext or ext not in ['.pdf', '.docx', '.txt', '.pptx', '.xlsx', '.xls', '.csv']:
            try:
                import magic
                mime = magic.Magic(mime=True)
                mime_type = mime.from_file(file_path)
                if mime_type == 'application/pdf':
                    ext = '.pdf'
                elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    ext = '.docx'
                elif mime_type == 'text/plain':
                    ext = '.txt'
                elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    ext = '.pptx'
                elif mime_type.startswith('application/vnd.ms-excel') or mime_type.startswith(
                        'application/vnd.openxmlformats-officedocument.spreadsheetml'):
                    ext = '.xlsx'
                elif mime_type == 'text/csv':
                    ext = '.csv'
                logger.info(f"Detected MIME type: {mime_type} -> {ext}")
            except ImportError:
                file_name = os.path.basename(file_path).lower()
                if 'pdf' in file_name:
                    ext = '.pdf'
                elif 'doc' in file_name:
                    ext = '.docx'
                elif 'txt' in file_name or 'text' in file_name:
                    ext = '.txt'
                elif 'ppt' in file_name:
                    ext = '.pptx'
                elif 'xl' in file_name or 'excel' in file_name:
                    ext = '.xlsx'
                elif 'csv' in file_name:
                    ext = '.csv'
                logger.info(f"Guessed file type from name: {ext}")
            except Exception as e:
                logger.warning(f"Failed to detect file type: {str(e)}")

        logger.info(f"Processing file: {file_path}, extension: {ext}")

        if ext == '.pdf':
            return self.parse_pdf(file_path)
        elif ext == '.docx':
            return self.parse_docx(file_path)
        elif ext == '.txt':
            return self.parse_txt(file_path)
        elif ext == '.pptx':
            return self.parse_pptx(file_path)
        elif ext in ['.xlsx', '.xls']:
            return self.parse_excel(file_path)
        elif ext in ['.csv']:
            return self.parse_csv(file_path)
        else:
            logger.warning(f"Unknown file format: {ext}, trying as plain text")
            try:
                return self.parse_txt(file_path)
            except Exception as txt_error:
                raise ValueError(f"Cannot process file: {file_path}, error: {str(txt_error)}")

    def parse_pdf(self, file_path):
        if not OPTIONAL_DEPS_AVAILABLE:
            raise ValueError("PDF parsing requires additional dependencies")

        doc = fitz.open(file_path)
        full_text = []

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            page_text = page.get_text()

            if len(page_text.strip()) < 50 and self.reader and not self.disable_ocr:
                pix = page.get_pixmap()
                img = Image.open(io.BytesIO(pix.tobytes("png")))
                img_np = np.array(img)
                ocr_results = self.reader.readtext(img_np)
                page_text_from_ocr = ' '.join([text for _, text, _ in ocr_results])
                if page_text_from_ocr.strip():
                    page_text = page_text_from_ocr

            full_text.append(page_text)

            if self.reader and not self.disable_ocr:
                image_list = page.get_images(full=True)
                for img_index, img_info in enumerate(image_list):
                    xref = img_info[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    try:
                        img = Image.open(io.BytesIO(image_bytes))
                        img_np = np.array(img)
                        ocr_results = self.reader.readtext(img_np)
                        image_text = ' '.join([text for _, text, _ in ocr_results])
                        if image_text.strip():
                            full_text.append(image_text)
                    except Exception as e:
                        logger.error(f"Error processing image: {e}")

        combined_text = '\n'.join(full_text)
        return self.segment_by_semantics(combined_text)

    def parse_docx(self, file_path):
        doc = docx.Document(file_path)
        paragraphs_text = []

        for para in doc.paragraphs:
            if para.text.strip():
                cleaned_text = self.clean_text(para.text)
                paragraphs_text.append(cleaned_text)

        table_texts = []
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                table_text.append(' | '.join(row_text))

            if table_text:
                table_content = '\n'.join(table_text)
                cleaned_table = self.clean_text(table_content)
                table_texts.append(cleaned_table)

        combined_text = '\n\n'.join(paragraphs_text + table_texts)
        return self.segment_by_semantics(combined_text)

    def parse_txt(self, file_path):
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
            text = file.read()
        return self.segment_by_semantics(text)

    def parse_pptx(self, file_path):
        if not OPTIONAL_DEPS_AVAILABLE:
            raise ValueError("PowerPoint parsing requires additional dependencies")

        prs = Presentation(file_path)
        slide_texts = []

        for slide in prs.slides:
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)

            if texts:
                slide_content = '\n'.join(texts)
                cleaned_slide = self.clean_text(slide_content)
                slide_texts.append(cleaned_slide)

        combined_text = '\n\n'.join(slide_texts)
        return self.segment_by_semantics(combined_text)

    def parse_excel(self, file_path):
        if not OPTIONAL_DEPS_AVAILABLE:
            raise ValueError("Excel parsing requires additional dependencies")

        wb = load_workbook(file_path, data_only=True)
        sheet_texts = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            rows = []

            for row in sheet.iter_rows(values_only=True):
                if any(cell for cell in row):
                    row_text = ' | '.join(str(cell) for cell in row if cell is not None)
                    rows.append(row_text)

            if rows:
                sheet_content = '\n'.join(rows)
                cleaned_sheet = self.clean_text(sheet_content)
                sheet_texts.append(cleaned_sheet)

        combined_text = '\n\n'.join(sheet_texts)
        return self.segment_by_semantics(combined_text)

    def parse_csv(self, file_path):
        rows = []

        with open(file_path, 'r', newline='', encoding='utf-8', errors='ignore') as csvfile:
            csv_reader = csv.reader(csvfile)
            for row in csv_reader:
                if any(cell.strip() for cell in row):
                    rows.append(' | '.join(row))

        if rows:
            combined_text = '\n'.join(rows)
            return self.segment_by_semantics(combined_text)
        return []

    def clean_text(self, text):
        if not text or not text.strip():
            return text

        if self.replace_whitespace:
            text = re.sub(r'\s+', ' ', text)

        if self.remove_urls_emails:
            text = re.sub(r'https?://\S+|www\.\S+', '', text)
            text = re.sub(r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', '', text)

        if self.clean_for_ai:
            text = self.clean_text_for_ai_training(text)

        return text.strip()

    def clean_text_for_ai_training(self, text):
        if not text or not text.strip():
            return text

        patterns = [
            (r'https?://\S+|www\.\S+', ''),
            (r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b', ''),
            (r'\b(?:\d{1,3}\.){3}\d{1,3}\b', ''),
            (r'\+?\d{1,4}?[-.\s]?\(?\d{1,3}?\)?[-.\s]?\d{1,4}[-.\s]?\d{1,4}[-.\s]?\d{1,9}', ''),
            (r'@\w+', ''),
            (r'\b(?:\d{4}[-\s]?){3}\d{4}\b', ''),
            (r'\b\d{3}-\d{2}-\d{4}\b', ''),
            (r'\b\d{17}[\dXx]\b', ''),
            (r'\b[A-Za-z0-9_\-]{20,}\b', ''),
            (r'\[.*?\]\(.*?\)', ''),
            (r'<[^>]+>', '')
        ]

        for pattern, replacement in patterns:
            text = re.sub(pattern, replacement, text)

        text = re.sub(r'\s+', ' ', text)

        while '  ' in text:
            text = text.replace('  ', ' ')

        return text.strip()

    def identify_semantic_blocks(self, text):
        if not text or not text.strip():
            return []

        text = re.sub(r'[ \t]+', ' ', text).strip()
        text = re.sub(r'\n{3,}', '\n\n', text)

        paragraph_blocks = re.split(r'\n\s*\n', text)
        semantic_blocks = []

        for block in paragraph_blocks:
            if not block.strip():
                continue

            if self.nlp and self.primary_language in ['en', 'zh']:
                doc = self.nlp(block)
                current_sentences = []
                sentence_count = 0

                for sent in doc.sents:
                    sentence_count += 1
                    current_sentences.append(sent.text)

                    if sentence_count >= 5:
                        semantic_blocks.append(current_sentences)
                        current_sentences = []
                        sentence_count = 0

                if current_sentences:
                    semantic_blocks.append(current_sentences)
            else:
                if self.primary_language == 'zh':
                    sentences = re.split(r'([。！？；])', block)
                    processed_sentences = []

                    i = 0
                    while i < len(sentences):
                        if i + 1 < len(sentences) and re.match(r'[。！？；]', sentences[i + 1]):
                            processed_sentences.append(sentences[i] + sentences[i + 1])
                            i += 2
                        else:
                            if sentences[i].strip():
                                processed_sentences.append(sentences[i])
                            i += 1
                else:
                    processed_sentences = sent_tokenize(block)

                current_block = []
                for sent in processed_sentences:
                    current_block.append(sent)
                    if len(current_block) >= 5:
                        semantic_blocks.append(current_block)
                        current_block = []

                if current_block:
                    semantic_blocks.append(current_block)

        return semantic_blocks

    def split_chinese_sentences(self, text):
        sentences = re.split(r'([。！？；])', text)
        processed_sentences = []

        i = 0
        while i < len(sentences):
            if i + 1 < len(sentences) and re.match(r'[。！？；]', sentences[i + 1]):
                processed_sentences.append(sentences[i] + sentences[i + 1])
                i += 2
            else:
                if sentences[i].strip():
                    processed_sentences.append(sentences[i])
                i += 1

        return processed_sentences

    def segment_by_semantics(self, text):
        if not text or not text.strip():
            return []

        text = self.clean_text(text)

        semantic_blocks = self.identify_semantic_blocks(text)
        segments = []
        last_part = None

        for block in semantic_blocks:
            block_text = ' '.join(block)

            if len(block_text) > self.segment_size:
                sub_segments = self.split_large_block(block)

                for segment in sub_segments:
                    if last_part and self.overlap:
                        segment = self.apply_overlap(segment, last_part)

                    segments.append({
                        'content': segment,
                        'type': 'semantic_unit'
                    })
                    last_part = segment
            else:
                if last_part and self.overlap:
                    block_text = self.apply_overlap(block_text, last_part)

                segments.append({
                    'content': block_text,
                    'type': 'semantic_unit'
                })
                last_part = block_text

        optimized_segments = self.optimize_segment_sizes(segments)
        return optimized_segments

    def split_large_block(self, sentences):
        segments = []
        current_segment = []
        current_length = 0

        for sentence in sentences:
            if len(sentence) > self.segment_size:
                if current_segment:
                    segments.append(' '.join(current_segment))
                    current_segment = []
                    current_length = 0

                split_parts = self.split_long_sentence(sentence)
                segments.extend(split_parts)
                continue

            if current_length + len(sentence) + 1 <= self.segment_size:
                current_segment.append(sentence)
                current_length += len(sentence) + 1
            else:
                if current_segment:
                    segments.append(' '.join(current_segment))

                current_segment = [sentence]
                current_length = len(sentence)

        if current_segment:
            segments.append(' '.join(current_segment))

        return segments

    def split_long_sentence(self, sentence):
        clause_markers = [
            ',', '，', ';', '；', ':', '：',
            'and', 'but', 'or', 'nor', 'yet', 'so',
            '和', '或者', '但是', '而且', '然后', '因此'
        ]

        parts = []
        current_part = ""
        words = re.split(r'(\s+)', sentence) if self.primary_language != 'zh' else list(sentence)

        for word in words:
            if current_part and len(current_part + word) > self.segment_size:
                parts.append(current_part.strip())
                current_part = word
            else:
                current_part += word

                for marker in clause_markers:
                    if word.endswith(marker) and len(current_part) > self.segment_size // 2:
                        parts.append(current_part.strip())
                        current_part = ""
                        break

        if current_part:
            parts.append(current_part.strip())

        final_parts = []
        for part in parts:
            if len(part) <= self.segment_size:
                final_parts.append(part)
            else:
                for i in range(0, len(part), self.segment_size):
                    final_parts.append(part[i:i + self.segment_size])

        return final_parts

    def optimize_segment_sizes(self, segments):
        if not segments or len(segments) < 2:
            return segments

        optimized = []
        i = 0

        while i < len(segments):
            current = segments[i]

            if (i + 1 < len(segments) and
                    len(current['content']) + len(segments[i + 1]['content']) + 1 <= self.segment_size):

                combined_content = current['content'] + ' ' + segments[i + 1]['content']
                optimized.append({
                    'content': combined_content,
                    'type': 'semantic_unit'
                })
                i += 2
            else:
                optimized.append(current)
                i += 1

        return optimized

    def apply_overlap(self, current_text, last_part):
        if not current_text.strip() or not self.overlap or not last_part:
            return current_text

        if self.nlp:
            doc = self.nlp(last_part)
            sentences = list(doc.sents)
            if sentences:
                last_sentence = sentences[-1].text
            else:
                return current_text
        else:
            if self.primary_language == 'zh':
                sentences = self.split_chinese_sentences(last_part)
            else:
                sentences = sent_tokenize(last_part)

            if sentences:
                last_sentence = sentences[-1]
            else:
                return current_text

        if len(last_sentence) > self.overlap_limit:
            last_sentence = last_sentence[-self.overlap_limit:]

        if current_text.startswith(last_sentence):
            return current_text

        return last_sentence + " " + current_text


def process_document(file_path, languages=None, segment_size=1000, overlap=True, overlap_limit=200, clean_for_ai=True,
                     replace_whitespace=False, remove_urls_emails=False,
                     disable_ocr=False):
    if languages is None:
        languages = ['ch_sim', 'en']

    parser = SemanticDocumentParser(
        language_list=languages,
        segment_size=segment_size,
        overlap=overlap,
        overlap_limit=overlap_limit,
        clean_for_ai=clean_for_ai,
        replace_whitespace=replace_whitespace,
        remove_urls_emails=remove_urls_emails,
        disable_ocr=disable_ocr
    )

    segments = parser.parse(file_path)
    return segments


def export_segments_to_xlsx(segments, output_file):
    if not OPTIONAL_DEPS_AVAILABLE:
        raise ValueError("Excel export requires openpyxl dependency")

    wb = Workbook()
    ws = wb.active
    ws.title = "Segments"

    header_font = Font(bold=True, size=12)
    header_fill = PatternFill(start_color="D9EAD3", end_color="D9EAD3", fill_type="solid")
    header_alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    border = Border(left=Side(style='thin'), right=Side(style='thin'), top=Side(style='thin'),
                    bottom=Side(style='thin'))

    headers = ["编号", "层级", "内容"]
    ws.append(headers)

    for col in range(1, 4):
        cell = ws.cell(row=1, column=col)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = header_alignment
        cell.border = border

    def detect_level(segment, index):
        if index == 0:
            content = segment.get('content', '')
            if len(content) < 100:
                return "标题1"

        content = segment.get('content', '')

        if re.match(r'^(第[一二三四五六七八九十]+章|第\d+章|Chapter \d+|I\.|II\.|III\.|IV\.|V\.)', content.strip()):
            return "标题1"

        if re.match(r'^(\d+\.\d+|\d+\.\d+\.\d+|[一二三四五六七八九十]+、|\(一\)|\(二\))', content.strip()):
            return "标题2"

        if re.match(r'^\d+\.\s|\(\d+\)\s|•\s|-\s', content.strip()):
            return "标题3"

        if len(content) < 100:
            return "可能标题"
        else:
            return "正文"

    for i, segment in enumerate(segments):
        level_name = detect_level(segment, i)
        content = segment.get('content', '')
        ws.append([i + 1, level_name, content])

    ws.column_dimensions['A'].width = 10
    ws.column_dimensions['B'].width = 25
    ws.column_dimensions['C'].width = 100

    for row in range(2, len(segments) + 2):
        level_cell = ws.cell(row=row, column=2).value

        if "标题1" in str(level_cell):
            ws.cell(row=row, column=2).fill = PatternFill(start_color="FFD966", end_color="FFD966", fill_type="solid")
            ws.cell(row=row, column=3).font = Font(bold=True, size=12)
        elif "标题2" in str(level_cell):
            ws.cell(row=row, column=2).fill = PatternFill(start_color="A9D08E", end_color="A9D08E", fill_type="solid")
            ws.cell(row=row, column=3).font = Font(bold=True)
        elif "标题3" in str(level_cell):
            ws.cell(row=row, column=2).fill = PatternFill(start_color="9BC2E6", end_color="9BC2E6", fill_type="solid")
            ws.cell(row=row, column=3).font = Font(italic=True)

    wb.save(output_file)
    return output_file


class ParallelDialogueGenerator:
    def __init__(self, api_url=None, api_key=None, output_dir="dialogues",
                 model="gpt-4o-mini", max_generation_mode=False,
                 max_questions=10, dialogue_token_limit=1000,
                 language="en-zh", file_identifier=None,
                 use_api=True, require_api_key=False):
        self.use_api = use_api
        self.require_api_key = require_api_key

        if api_url:
            self.api_base = api_url.rstrip('/')
            if not self.api_base.endswith('/v1'):
                if '/v1' not in self.api_base:
                    self.api_base += '/v1'
            self.api_url = f"{self.api_base}/chat/completions"
        else:
            self.api_base = None
            self.api_url = None

        self.api_key = api_key
        self.model = model
        self.output_dir = output_dir
        self.language = language
        self.file_identifier = file_identifier

        self.max_generation_mode = max_generation_mode
        self.max_questions = max_questions
        self.dialogue_token_limit = dialogue_token_limit

        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def _call_api(self, messages, max_tokens=None, max_retries=3):
        if not self.use_api or not self.api_url:
            prompt = ""
            for msg in messages:
                if msg['role'] == 'user':
                    prompt = msg['content']
                    break

            placeholder = self._generate_placeholder_response(prompt)
            return {
                "choices": [
                    {
                        "message": {
                            "content": placeholder,
                            "role": "assistant"
                        }
                    }
                ]
            }

        headers = {
            "Content-Type": "application/json; charset=utf-8",
        }

        if self.api_key:
            headers["Authorization"] = f"Bearer {self.api_key}"

        max_tokens = max_tokens or (self.dialogue_token_limit if self.max_generation_mode else 1000)

        data = {
            "model": self.model,
            "messages": messages,
            "max_tokens": max_tokens,
            "temperature": 0.7
        }

        for attempt in range(max_retries):
            try:
                logger.info(f"API call attempt {attempt + 1}")
                json_data = json.dumps(data, ensure_ascii=False).encode('utf-8')

                response = requests.post(
                    self.api_url,
                    headers=headers,
                    data=json_data,
                    timeout=30
                )

                logger.info(f"Response status: {response.status_code}")
                if response.status_code != 200:
                    logger.error(f"Error: {response.text[:200]}...")

                response.raise_for_status()
                result = response.json()

                if "choices" not in result:
                    raise ValueError("Invalid API response format")

                return result
            except Exception as e:
                if attempt == max_retries - 1:
                    raise Exception(f"API call failed after {max_retries} attempts: {str(e)}")
                logger.error(f"API call failed: {str(e)}. Retrying...")
                time.sleep(2 ** attempt)

    def _generate_placeholder_response(self, prompt):
        if "generate questions" in prompt.lower():
            if self.language == "zh":
                return "1. 这个文本的主要论点是什么？\n2. 作者使用了哪些证据来支持他们的观点？\n3. 文中提到的关键概念如何应用于实际情况？"
            elif self.language == "en":
                return "1. What is the main argument of this text?\n2. What evidence does the author use to support their claims?\n3. How do the key concepts mentioned apply to real-world situations?"
            else:
                return "1. What is the main argument of this text? 这个文本的主要论点是什么？\n2. What evidence does the author use to support their claims? 作者使用了哪些证据来支持他们的观点？\n3. How do the key concepts mentioned apply to real-world situations? 文中提到的关键概念如何应用于实际情况？"

        elif "follow-up question" in prompt.lower():
            if self.language == "zh":
                return "您能进一步解释这个观点对当今社会的影响吗？"
            elif self.language == "en":
                return "Could you elaborate on the impact this perspective has on today's society?"
            else:
                return "Could you elaborate on the impact this perspective has on today's society? 您能进一步解释这个观点对当今社会的影响吗？"

        else:
            if self.language == "zh":
                return "这是一个自动生成的回复。在没有API连接的情况下，系统无法生成真实的对话内容。这是一个示例回复，用于演示界面功能。在实际使用中，您需要提供有效的API密钥来获取真实的AI生成内容。"
            elif self.language == "en":
                return "This is an automatically generated response. Without an API connection, the system cannot generate real dialogue content. This is a sample response to demonstrate the interface functionality. In actual use, you need to provide a valid API key to get real AI-generated content."
            else:
                return "This is an automatically generated response. Without an API connection, the system cannot generate real dialogue content.\n\n这是一个自动生成的回复。在没有API连接的情况下，系统无法生成真实的对话内容。"

    def generate_questions(self, chunk, num_questions=None):
        if num_questions is None:
            num_questions = self.max_questions if self.max_generation_mode else 3

        if self.language == "en":
            language_instruction = "Generate questions in English only."
        elif self.language == "zh":
            language_instruction = "请仅用中文生成问题。"
        else:
            language_instruction = "Generate each question in both English and Chinese (同时用英文和中文生成每个问题)."

        prompt = f"""Please generate {num_questions} insightful questions based on the following text.
        Focus on key points, implications, and potential applications of the information.
        {language_instruction}

        {chunk}

        Return only questions in a numbered list form, without any additional text."""

        messages = [
            {"role": "system",
             "content": "You are an assistant who generates insightful questions in both English and Chinese."},
            {"role": "user", "content": prompt}
        ]

        max_tokens = 2000 if self.max_generation_mode else 1000
        response = self._call_api(messages, max_tokens=max_tokens)

        questions_text = response.get("choices", [{}])[0].get("message", {}).get("content", "")
        questions = []

        for line in questions_text.strip().split('\n'):
            line = line.strip()
            if not line:
                continue

            if (line[0].isdigit() or
                    line.startswith('- ') or
                    line.startswith('• ') or
                    line.startswith('* ')):

                if '. ' in line and line[0].isdigit():
                    question = line.split('. ', 1)[-1]
                elif '、' in line and line[0].isdigit():
                    question = line.split('、', 1)[-1]
                elif line.startswith(('- ', '• ', '* ')):
                    question = line[2:]
                else:
                    question = line

                questions.append(question)
            elif questions:
                questions[-1] += " " + line

        return questions[:num_questions]

    def conduct_dialogue(self, question, chunk, rounds=3):
        if self.language == "en":
            system_content = f"You are having a conversation about the following text in English: {chunk}"
        elif self.language == "zh":
            system_content = f"你正在用中文讨论以下文本内容: {chunk}"
        else:
            system_content = f"You are having a bilingual conversation about the following text. Please reply in the same language as the user's question: {chunk}"

        messages = [
            {"role": "system", "content": system_content},
            {"role": "user", "content": question}
        ]

        conversation_history = messages.copy()

        for i in range(rounds):
            print(f"Dialogue round {i + 1}/{rounds} for question: {question[:50]}...")

            response = self._call_api(conversation_history)
            assistant_message = {
                "role": "assistant",
                "content": response.get("choices", [{}])[0].get("message", {}).get("content", "")
            }
            conversation_history.append(assistant_message)

            if i < rounds - 1:
                if self.language == "en":
                    follow_up_instruction = "Based on the previous exchange, generate a natural follow-up question in English to deepen the conversation."
                elif self.language == "zh":
                    follow_up_instruction = "根据前面的对话，用中文生成一个自然的后续问题，以深入讨论。"
                else:
                    follow_up_instruction = "Based on the previous exchange, generate a natural follow-up question to deepen the conversation. Use the same language (English or Chinese) as the most recent reply."

                follow_up_messages = [
                    {"role": "system", "content": follow_up_instruction},
                    {"role": "user",
                     "content": f"Generate a follow-up question for this response: {assistant_message['content']}"}
                ]
                follow_up_response = self._call_api(follow_up_messages)
                follow_up_question = follow_up_response.get("choices", [{}])[0].get("message", {}).get("content", "")

                conversation_history.append({"role": "user", "content": follow_up_question})

        return conversation_history

    def process_chunk(self, chunk, chunk_id, rounds=3):
        try:
            print(f"Processing text chunk {chunk_id}...")

            questions = self.generate_questions(chunk)
            print(f"Generated {len(questions)} questions for chunk {chunk_id}")

            dialogues = {}
            for i, question in enumerate(questions):
                print(f"Conducting dialogue {i + 1}/{len(questions)} for chunk {chunk_id}")
                dialogue = self.conduct_dialogue(question, chunk, rounds)
                dialogues[f"question_{i + 1}"] = {
                    "question": question,
                    "dialogue": dialogue
                }

            result = {
                "chunk_id": chunk_id,
                "chunk_text": chunk,
                "questions": questions,
                "dialogues": dialogues
            }

            self._save_dialogue(result, chunk_id)
            print(f"Successfully completed processing for chunk {chunk_id}")

            return result
        except Exception as e:
            print(f"Error processing chunk {chunk_id}: {str(e)}")
            return {
                "chunk_id": chunk_id,
                "error": str(e),
                "status": "failed"
            }

    def process_chunks_parallel(self, chunks, rounds=3, max_workers=4):
        results = []

        with concurrent.futures.ThreadPoolExecutor(max_workers=max_workers) as executor:
            future_to_chunk = {
                executor.submit(self.process_chunk, chunk, i, rounds): (i, chunk)
                for i, chunk in enumerate(chunks)
            }

            for future in concurrent.futures.as_completed(future_to_chunk):
                chunk_id, chunk = future_to_chunk[future]
                try:
                    result = future.result()
                    results.append(result)
                    if "status" in result and result["status"] == "failed":
                        print(f"Failed processing chunk {chunk_id}: {result.get('error', 'Unknown error')}")
                    else:
                        print(f"Completed processing for chunk {chunk_id}")
                except Exception as e:
                    print(f"Error processing chunk {chunk_id}: {str(e)}")
                    results.append({
                        "chunk_id": chunk_id,
                        "error": str(e),
                        "status": "failed"
                    })

        results.sort(key=lambda x: x.get("chunk_id", 0))
        return results

    def _save_dialogue(self, dialogue_data, chunk_id):
        filename = f"dialogue_chunk_{chunk_id}.json"
        filepath = os.path.join(self.output_dir, filename)

        with open(filepath, 'w', encoding='utf-8') as f:
            json.dump(dialogue_data, f, ensure_ascii=False, indent=2)

        return filepath


def process_document_and_generate_dialogues(
        file_path,
        api_url=None,
        api_key=None,
        model="gpt-4o-mini",
        dialogue_rounds=3,
        max_workers=8,
        output_dir="results",
        language="en-zh",
        max_generation_mode=False,
        max_questions=5,
        dialogue_token_limit=1000,
        segment_size=1000,
        overlap=True,
        overlap_limit=200,
        clean_for_ai=True,
        replace_whitespace=False,
        remove_urls_emails=False,
        languages=None,
        file_identifier=None,
        use_api=True,
        require_api_key=False,
        disable_ocr=False
):
    if languages is None:
        languages = ['ch_sim', 'en']

    output_dir = os.path.abspath(output_dir)
    if not os.path.exists(output_dir):
        os.makedirs(output_dir)

    dialogue_dir = os.path.join(output_dir, "dialogues")
    if not os.path.exists(dialogue_dir):
        os.makedirs(dialogue_dir)

    logger.info(f"Processing document: {file_path}")
    try:
        segments = process_document(
            file_path,
            languages=languages,
            segment_size=segment_size,
            overlap=overlap,
            overlap_limit=overlap_limit,
            clean_for_ai=clean_for_ai,
            replace_whitespace=replace_whitespace,
            remove_urls_emails=remove_urls_emails,
            disable_ocr=disable_ocr
        )

        logger.info(f"Document processed into {len(segments)} segments")

        segments_file = os.path.join(output_dir, "segments.json")
        with open(segments_file, 'w', encoding='utf-8') as f:
            json.dump(segments, f, ensure_ascii=False, indent=2)

        xlsx_file = None
        if OPTIONAL_DEPS_AVAILABLE:
            xlsx_file = os.path.join(output_dir, "segments.xlsx")
            export_segments_to_xlsx(segments, xlsx_file)
            logger.info(f"Segments exported to Excel: {xlsx_file}")

        chunks = [segment['content'] for segment in segments]

        generator = ParallelDialogueGenerator(
            api_url=api_url,
            api_key=api_key,
            output_dir=dialogue_dir,
            model=model,
            max_generation_mode=max_generation_mode,
            max_questions=max_questions,
            dialogue_token_limit=dialogue_token_limit,
            language=language,
            file_identifier=file_identifier,
            use_api=use_api,
            require_api_key=require_api_key
        )

        logger.info(f"Generating dialogues for {len(chunks)} segments")
        dialogue_results = generator.process_chunks_parallel(
            chunks=chunks,
            rounds=dialogue_rounds,
            max_workers=max_workers
        )

        combined_results = {
            "document": os.path.basename(file_path),
            "file_identifier": file_identifier,
            "segments": segments,
            "dialogues": dialogue_results
        }

        results_file = os.path.join(output_dir, "combined_results.json")
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(combined_results, f, ensure_ascii=False, indent=2)

        return {
            "status": "success",
            "segments_count": len(segments),
            "dialogues_count": len(dialogue_results),
            "results_file": results_file,
            "xlsx_file": xlsx_file
        }

    except Exception as e:
        logger.error(f"Error in integrated workflow: {str(e)}")
        return {
            "status": "error",
            "error": str(e)
        }


class HistoryManager:
    def __init__(self, history_dir='history'):
        self.history_dir = history_dir
        if not os.path.exists(history_dir):
            os.makedirs(history_dir)

        self.history_file = os.path.join(history_dir, 'processing_history.json')
        self._init_history_file()

    def _init_history_file(self):
        if not os.path.exists(self.history_file):
            with open(self.history_file, 'w', encoding='utf-8') as f:
                json.dump([], f, ensure_ascii=False)

    def get_all_records(self):
        try:
            with open(self.history_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, FileNotFoundError) as e:
            logger.error(f"Error reading history file: {e}")
            return []

    def get_record(self, job_id):
        records = self.get_all_records()
        for record in records:
            if record.get('job_id') == job_id:
                return record
        return None

    def add_record(self, record):
        if 'job_id' not in record:
            record['job_id'] = str(uuid.uuid4())

        if 'timestamp' not in record:
            record['timestamp'] = datetime.now().strftime('%Y-%m-%d %H:%M:%S')

        records = self.get_all_records()
        records.append(record)

        if len(records) > 100:
            records = records[-100:]

        try:
            with open(self.history_file, 'w', encoding='utf-8') as f:
                json.dump(records, f, ensure_ascii=False, indent=2)
            return record
        except Exception as e:
            logger.error(f"Error saving history record: {e}")
            return None

    def clear_history(self):
        try:
            with open(self.history_file, 'w', encoding='utf-8') as f:
                json.dump([], f, ensure_ascii=False)
            return True
        except Exception as e:
            logger.error(f"Error clearing history: {e}")
            return False


app = Flask(__name__)
CORS(app)

MAX_CONTENT_LENGTH = int(os.environ.get('MAX_CONTENT_LENGTH', 50 * 1024 * 1024))
app.config['MAX_CONTENT_LENGTH'] = MAX_CONTENT_LENGTH

UPLOAD_FOLDER = os.environ.get('UPLOAD_FOLDER')
upload_paths = [
    UPLOAD_FOLDER,
    '/tmp/uploads',
    '/var/tmp/uploads',
    os.path.join(os.getcwd(), 'uploads'),
    tempfile.gettempdir()
]

for path in upload_paths:
    if path and ensure_directory(path):
        UPLOAD_FOLDER = path
        logger.info(f"Using upload folder: {UPLOAD_FOLDER}")
        break
else:
    UPLOAD_FOLDER = tempfile.mkdtemp()
    logger.warning(f"Falling back to temporary directory: {UPLOAD_FOLDER}")

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER

RESULTS_FOLDER = os.environ.get('RESULTS_FOLDER', 'results')
if not os.path.exists(RESULTS_FOLDER):
    os.makedirs(RESULTS_FOLDER)
app.config['RESULTS_FOLDER'] = RESULTS_FOLDER

HISTORY_FOLDER = os.environ.get('HISTORY_FOLDER', 'history')
if not os.path.exists(HISTORY_FOLDER):
    os.makedirs(HISTORY_FOLDER)
app.config['HISTORY_FOLDER'] = HISTORY_FOLDER

history_manager = HistoryManager(HISTORY_FOLDER)


def handle_file_upload(request_files, upload_folder):
    file_paths = []
    filenames = []

    for file in request_files:
        filename = str(uuid.uuid4()) + '_' + secure_filename(file.filename)
        file_path = os.path.join(upload_folder, filename)
        file.save(file_path)
        file_paths.append(file_path)
        filenames.append(file.filename)

    return file_paths, filenames


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/api/history', methods=['GET'])
def api_get_history():
    try:
        records = history_manager.get_all_records()

        limit = request.args.get('limit')
        job_type = request.args.get('job_type')
        status = request.args.get('status')

        if job_type:
            records = [r for r in records if r.get('job_type') == job_type]

        if status:
            records = [r for r in records if r.get('status') == status]

        records = sorted(records, key=lambda r: r.get('timestamp', ''), reverse=True)

        if limit:
            try:
                limit = int(limit)
                records = records[:limit]
            except ValueError:
                pass

        return jsonify({
            "status": "success",
            "count": len(records),
            "records": records
        })
    except Exception as e:
        logger.error(f"Error retrieving history: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/history/<job_id>', methods=['GET'])
def api_get_history_record(job_id):
    try:
        record = history_manager.get_record(job_id)
        if record:
            return jsonify({
                "status": "success",
                "record": record
            })
        else:
            return jsonify({"error": "Record not found"}), 404
    except Exception as e:
        logger.error(f"Error retrieving history record: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/history/clear', methods=['POST'])
def api_clear_history():
    try:
        success = history_manager.clear_history()
        if success:
            return jsonify({"status": "success", "message": "History cleared"})
        else:
            return jsonify({"error": "Failed to clear history"}), 500
    except Exception as e:
        logger.error(f"Error clearing history: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/process-document', methods=['POST'])
def api_process_document():
    logger.info("Received document processing request")

    job_start_time = time.time()
    job_id = str(uuid.uuid4())

    if 'files' not in request.files:
        return jsonify({"error": "No files part in the request"}), 400

    files = request.files.getlist('files')
    if not files or files[0].filename == '':
        return jsonify({"error": "No files selected"}), 400

    try:
        languages = request.form.get('languages', 'ch_sim,en').split(',')
        segment_size = int(request.form.get('segment_size', 1000))
        overlap = request.form.get('overlap', 'true').lower() == 'true'
        overlap_limit = int(request.form.get('overlap_limit', 200))
        clean_for_ai = request.form.get('clean_for_ai', 'true').lower() == 'true'
        replace_whitespace = request.form.get('replace_whitespace', 'false').lower() == 'true'
        remove_urls_emails = request.form.get('remove_urls_emails', 'false').lower() == 'true'
        disable_ocr = request.form.get('disable_ocr', 'false').lower() == 'true'
    except Exception as e:
        history_manager.add_record({
            "job_id": job_id,
            "job_type": "document_processing",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "error",
            "error": f"Invalid parameters: {str(e)}",
            "files": [f.filename for f in files],
            "duration_seconds": time.time() - job_start_time
        })
        return jsonify({"error": f"Invalid parameters: {str(e)}"}), 400

    try:
        results = {}
        total_segments = 0
        export_format = request.form.get('export_format', 'json').lower()
        result_files = {}

        result_dir = os.path.join(app.config['RESULTS_FOLDER'], str(uuid.uuid4()))
        os.makedirs(result_dir)

        processed_files = []
        for file in files:
            try:
                filename = str(uuid.uuid4()) + '_' + secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(file_path)
                processed_files.append(file.filename)

                file_result_dir = os.path.join(result_dir, secure_filename(file.filename))
                os.makedirs(file_result_dir)

                segments = process_document(
                    file_path,
                    languages=languages,
                    segment_size=segment_size,
                    overlap=overlap,
                    overlap_limit=overlap_limit,
                    clean_for_ai=clean_for_ai,
                    replace_whitespace=replace_whitespace,
                    remove_urls_emails=remove_urls_emails,
                    disable_ocr=disable_ocr
                )

                total_segments += len(segments)

                json_file = os.path.join(file_result_dir, "segments.json")
                with open(json_file, 'w', encoding='utf-8') as f:
                    json.dump(segments, f, ensure_ascii=False, indent=2)

                xlsx_file = None
                if OPTIONAL_DEPS_AVAILABLE:
                    xlsx_file = os.path.join(file_result_dir, "segments.xlsx")
                    export_segments_to_xlsx(segments, xlsx_file)

                try:
                    os.remove(file_path)
                except Exception as e:
                    logger.warning(f"Failed to remove temp file {file_path}: {e}")

                results[file.filename] = {
                    "status": "success",
                    "segments_count": len(segments),
                    "segments": segments
                }

                result_files[file.filename] = {
                    "json": json_file,
                    "xlsx": xlsx_file
                }

            except Exception as e:
                logger.error(f"Error processing file {file.filename}: {str(e)}")
                results[file.filename] = {
                    "status": "error",
                    "error": str(e)
                }

        history_record = {
            "job_id": job_id,
            "job_type": "document_processing",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "success",
            "files": processed_files,
            "total_segments": total_segments,
            "parameters": {
                "languages": languages,
                "segment_size": segment_size,
                "overlap": overlap,
                "overlap_limit": overlap_limit,
                "clean_for_ai": clean_for_ai,
                "replace_whitespace": replace_whitespace,
                "remove_urls_emails": remove_urls_emails,
                "disable_ocr": disable_ocr
            },
            "result_directory": result_dir,
            "duration_seconds": time.time() - job_start_time
        }
        history_manager.add_record(history_record)

        return jsonify({
            "status": "success",
            "job_id": job_id,
            "processed_files": len(files),
            "total_segments": total_segments,
            "results": results,
            "result_files": result_files,
            "export_format": export_format
        })

    except Exception as e:
        history_manager.add_record({
            "job_id": job_id,
            "job_type": "document_processing",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "error",
            "error": str(e),
            "files": [f.filename for f in files],
            "duration_seconds": time.time() - job_start_time
        })
        logger.error(f"Error processing document: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/generate-dialogues', methods=['POST'])
def api_generate_dialogues():
    logger.info("Received dialogue generation request")

    job_start_time = time.time()
    job_id = str(uuid.uuid4())

    if not request.json or 'segments' not in request.json:
        return jsonify({"error": "No segments provided"}), 400

    try:
        segments = request.json['segments']
        api_url = request.json.get('api_url')
        api_key = request.json.get('api_key')
        use_api = request.json.get('use_api', True)
        require_api_key = request.json.get('require_api_key', False)

        if use_api and not api_url:
            return jsonify({"error": "API URL required when using API mode"}), 400

        if use_api and require_api_key and not api_key:
            return jsonify({"error": "API key required"}), 400

        model = request.json.get('model', 'gpt-4o-mini')
        dialogue_rounds = int(request.json.get('dialogue_rounds', 3))
        max_workers = int(request.json.get('max_workers', 4))
        language = request.json.get('language', 'en-zh')
        max_generation_mode = request.json.get('max_generation_mode', False)
        max_questions = int(request.json.get('max_questions', 5))
        dialogue_token_limit = int(request.json.get('dialogue_token_limit', 1000))
    except Exception as e:
        history_manager.add_record({
            "job_id": job_id,
            "job_type": "dialogue_generation",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "error",
            "error": f"Invalid parameters: {str(e)}",
            "segments_count": len(request.json.get('segments', [])),
            "duration_seconds": time.time() - job_start_time
        })
        return jsonify({"error": f"Invalid parameters: {str(e)}"}), 400

    result_dir = os.path.join(app.config['RESULTS_FOLDER'], str(uuid.uuid4()))
    os.makedirs(result_dir)

    chunks = [segment['content'] for segment in segments]

    try:
        generator = ParallelDialogueGenerator(
            api_url=api_url,
            api_key=api_key,
            output_dir=result_dir,
            model=model,
            max_generation_mode=max_generation_mode,
            max_questions=max_questions,
            dialogue_token_limit=dialogue_token_limit,
            language=language,
            use_api=use_api,
            require_api_key=require_api_key
        )

        dialogue_results = generator.process_chunks_parallel(
            chunks=chunks,
            rounds=dialogue_rounds,
            max_workers=max_workers
        )

        combined_results = {
            "segments": segments,
            "dialogues": dialogue_results
        }

        results_file = os.path.join(result_dir, "results.json")
        with open(results_file, 'w', encoding='utf-8') as f:
            json.dump(combined_results, f, ensure_ascii=False, indent=2)

        history_record = {
            "job_id": job_id,
            "job_type": "dialogue_generation",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "success",
            "segments_count": len(segments),
            "dialogues_count": len(dialogue_results),
            "parameters": {
                "model": model,
                "dialogue_rounds": dialogue_rounds,
                "language": language,
                "max_generation_mode": max_generation_mode,
                "max_questions": max_questions
            },
            "result_directory": result_dir,
            "results_file": results_file,
            "duration_seconds": time.time() - job_start_time
        }
        history_manager.add_record(history_record)

        return jsonify({
            "status": "success",
            "job_id": job_id,
            "dialogues_count": len(dialogue_results),
            "dialogues": dialogue_results,
            "results_file": results_file
        })

    except Exception as e:
        history_manager.add_record({
            "job_id": job_id,
            "job_type": "dialogue_generation",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "error",
            "error": str(e),
            "segments_count": len(segments),
            "duration_seconds": time.time() - job_start_time
        })
        logger.error(f"Error generating dialogues: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/integrated-workflow', methods=['POST'])
def api_integrated_workflow():
    logger.info("Received integrated workflow request")

    job_start_time = time.time()
    job_id = str(uuid.uuid4())

    if 'files' not in request.files:
        return jsonify({"error": "No files part in the request"}), 400

    files = request.files.getlist('files')
    if not files or files[0].filename == '':
        return jsonify({"error": "No files selected"}), 400

    use_api = request.form.get('use_api', 'true').lower() == 'true'
    require_api_key = request.form.get('require_api_key', 'false').lower() == 'true'
    api_url = request.form.get('api_url')
    api_key = request.form.get('api_key')

    if use_api and not api_url:
        return jsonify({"error": "API URL required when using API mode"}), 400

    if use_api and require_api_key and not api_key:
        return jsonify({"error": "API key required"}), 400

    try:
        languages = request.form.get('languages', 'ch_sim,en').split(',')
        segment_size = int(request.form.get('segment_size', 1000))
        overlap = request.form.get('overlap', 'true').lower() == 'true'
        overlap_limit = int(request.form.get('overlap_limit', 200))
        clean_for_ai = request.form.get('clean_for_ai', 'true').lower() == 'true'
        replace_whitespace = request.form.get('replace_whitespace', 'false').lower() == 'true'
        remove_urls_emails = request.form.get('remove_urls_emails', 'false').lower() == 'true'
        disable_ocr = request.form.get('disable_ocr', 'false').lower() == 'true'

        model = request.form.get('model', 'gpt-4o-mini')
        dialogue_rounds = int(request.form.get('dialogue_rounds', 3))
        max_workers = int(request.form.get('max_workers', 4))
        language = request.form.get('language', 'en-zh')
        max_generation_mode = request.form.get('max_generation_mode', 'false').lower() == 'true'
        max_questions = int(request.form.get('max_questions', 5))
        dialogue_token_limit = int(request.form.get('dialogue_token_limit', 1000))
    except Exception as e:
        history_manager.add_record({
            "job_id": job_id,
            "job_type": "integrated_workflow",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "error",
            "error": f"Invalid parameters: {str(e)}",
            "files": [f.filename for f in files],
            "duration_seconds": time.time() - job_start_time
        })
        return jsonify({"error": f"Invalid parameters: {str(e)}"}), 400

    result_dir = os.path.join(app.config['RESULTS_FOLDER'], str(uuid.uuid4()))
    os.makedirs(result_dir)

    try:
        all_results = []
        file_results = {}
        processed_files = []

        for file in files:
            try:
                filename = str(uuid.uuid4()) + '_' + secure_filename(file.filename)
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                file.save(file_path)
                processed_files.append(file.filename)

                file_result_dir = os.path.join(result_dir, secure_filename(file.filename))
                os.makedirs(file_result_dir)

                result = process_document_and_generate_dialogues(
                    file_path=file_path,
                    api_url=api_url,
                    api_key=api_key,
                    model=model,
                    dialogue_rounds=dialogue_rounds,
                    max_workers=max_workers,
                    output_dir=file_result_dir,
                    language=language,
                    max_generation_mode=max_generation_mode,
                    max_questions=max_questions,
                    dialogue_token_limit=dialogue_token_limit,
                    segment_size=segment_size,
                    overlap=overlap,
                    overlap_limit=overlap_limit,
                    clean_for_ai=clean_for_ai,
                    replace_whitespace=replace_whitespace,
                    remove_urls_emails=remove_urls_emails,
                    languages=languages,
                    file_identifier=file.filename,
                    use_api=use_api,
                    require_api_key=require_api_key,
                    disable_ocr=disable_ocr
                )

                try:
                    os.remove(file_path)
                except Exception as e:
                    logger.warning(f"Failed to remove temp file {file_path}: {e}")

                if result["status"] == "success":
                    with open(result["results_file"], 'r', encoding='utf-8') as f:
                        combined_results = json.load(f)

                    file_results[file.filename] = {
                        "status": "success",
                        "segments_count": result["segments_count"],
                        "dialogues_count": result["dialogues_count"],
                        "results": combined_results,
                        "results_file": result["results_file"],
                        "xlsx_file": result.get("xlsx_file")
                    }
                    all_results.append(file_results[file.filename])
                else:
                    file_results[file.filename] = {
                        "status": "error",
                        "error": result["error"]
                    }
                    all_results.append(file_results[file.filename])

            except Exception as e:
                logger.error(f"Error processing file {file.filename}: {str(e)}")
                file_results[file.filename] = {
                    "status": "error",
                    "error": str(e)
                }
                all_results.append(file_results[file.filename])

        combined_results_path = os.path.join(result_dir, "all_files_results.json")
        with open(combined_results_path, 'w', encoding='utf-8') as f:
            json.dump(file_results, f, ensure_ascii=False, indent=2)

        successful_files = sum(1 for r in all_results if r.get("status") == "success")

        history_record = {
            "job_id": job_id,
            "job_type": "integrated_workflow",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "success" if successful_files > 0 else "partial_success",
            "files": processed_files,
            "successful_files": successful_files,
            "failed_files": len(files) - successful_files,
            "parameters": {
                "document": {
                    "languages": languages,
                    "segment_size": segment_size,
                    "overlap": overlap,
                    "clean_for_ai": clean_for_ai
                },
                "dialogue": {
                    "model": model,
                    "dialogue_rounds": dialogue_rounds,
                    "language": language,
                    "max_questions": max_questions
                }
            },
            "result_directory": result_dir,
            "combined_results_file": combined_results_path,
            "duration_seconds": time.time() - job_start_time
        }
        history_manager.add_record(history_record)

        return jsonify({
            "status": "success",
            "job_id": job_id,
            "processed_files": len(files),
            "successful_files": successful_files,
            "results": file_results,
            "combined_results_file": combined_results_path
        })

    except Exception as e:
        history_manager.add_record({
            "job_id": job_id,
            "job_type": "integrated_workflow",
            "timestamp": datetime.now().strftime('%Y-%m-%d %H:%M:%S'),
            "status": "error",
            "error": str(e),
            "files": [f.filename for f in files],
            "duration_seconds": time.time() - job_start_time
        })
        logger.error(f"Error in integrated workflow: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/download-segments/<format>/<path:filename>', methods=['GET'])
def download_segments(format, filename):
    if format not in ['json', 'xlsx']:
        return jsonify({"error": "Invalid format. Use 'json' or 'xlsx'."}), 400

    try:
        result_dir = os.path.join(app.config['RESULTS_FOLDER'], filename)
        file_path = os.path.join(result_dir, f"segments.{format}")

        if not os.path.exists(file_path):
            return jsonify({"error": f"File not found: {file_path}"}), 404

        return send_file(
            file_path,
            as_attachment=True,
            download_name=f"segments.{format}",
            mimetype='application/json' if format == 'json' else 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
    except Exception as e:
        logger.error(f"Error downloading file: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/download-results', methods=['POST'])
def api_download_results():
    logger.info("Received download results request")

    if not request.json:
        return jsonify({"error": "No JSON data provided"}), 400

    try:
        fd, temp_file_path = tempfile.mkstemp(suffix='.json')
        with os.fdopen(fd, 'w', encoding='utf-8') as f:
            json.dump(request.json, f, ensure_ascii=False, indent=2)

        return send_file(
            temp_file_path,
            as_attachment=True,
            download_name="document_dialogue_results.json",
            mimetype='application/json'
        )
    except Exception as e:
        logger.error(f"Error creating download file: {e}")
        return jsonify({"error": f"Error creating download file: {str(e)}"}), 500


@app.route('/api/download-xlsx', methods=['POST'])
def api_download_xlsx():
    logger.info("Received download XLSX request")

    if not request.json or 'segments' not in request.json:
        return jsonify({"error": "No segments data provided"}), 400

    try:
        segments = request.json['segments']
        fd, temp_file_path = tempfile.mkstemp(suffix='.xlsx')
        os.close(fd)

        export_segments_to_xlsx(segments, temp_file_path)

        return send_file(
            temp_file_path,
            as_attachment=True,
            download_name="segments.xlsx",
            mimetype='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
    except Exception as e:
        logger.error(f"Error creating XLSX file: {e}")
        return jsonify({"error": f"Error creating XLSX file: {str(e)}"}), 500


@app.route('/api/history/<job_id>/results', methods=['GET'])
def api_get_job_results(job_id):
    try:
        record = history_manager.get_record(job_id)
        if not record:
            return jsonify({"error": "Job not found"}), 404

        results_file = record.get('results_file') or record.get('combined_results_file')

        if not results_file or not os.path.exists(results_file):
            result_dir = record.get('result_directory')
            if result_dir and os.path.exists(result_dir):
                for filename in ['results.json', 'combined_results.json', 'all_files_results.json']:
                    potential_file = os.path.join(result_dir, filename)
                    if os.path.exists(potential_file):
                        results_file = potential_file
                        break

        if not results_file or not os.path.exists(results_file):
            return jsonify({
                "status": "error",
                "message": "Results file not found"
            }), 404

        with open(results_file, 'r', encoding='utf-8') as f:
            results = json.load(f)

        return jsonify({
            "status": "success",
            "job_id": job_id,
            "job_type": record.get('job_type'),
            "results": results
        })

    except Exception as e:
        logger.error(f"Error retrieving job results: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/history/download/<job_id>', methods=['GET'])
def api_download_job_results(job_id):
    try:
        record = history_manager.get_record(job_id)
        if not record:
            return jsonify({"error": "Job not found"}), 404

        job_type = record.get('job_type', '')
        results_file = None

        for field in ['results_file', 'combined_results_file']:
            if field in record and record[field]:
                results_file = record[field]
                if os.path.exists(results_file):
                    break

        if not results_file or not os.path.exists(results_file):
            result_dir = record.get('result_directory')
            if result_dir and os.path.exists(result_dir):
                if job_type == 'document_processing':
                    potential_files = ['segments.json']
                elif job_type == 'dialogue_generation':
                    potential_files = ['results.json', 'all_dialogues.json']
                else:
                    potential_files = ['combined_results.json', 'all_files_results.json', 'results.json']

                for filename in potential_files:
                    potential_file = os.path.join(result_dir, filename)
                    if os.path.exists(potential_file):
                        results_file = potential_file
                        break

                if not results_file or not os.path.exists(results_file):
                    for root, dirs, files in os.walk(result_dir):
                        for file in files:
                            if file.endswith('.json'):
                                results_file = os.path.join(root, file)
                                break
                        if results_file and os.path.exists(results_file):
                            break

        if not results_file or not os.path.exists(results_file):
            return jsonify({
                "status": "error",
                "message": "Results file not found"
            }), 404

        if job_type == 'document_processing':
            download_name = f"document_segments_{job_id}.json"
        elif job_type == 'dialogue_generation':
            download_name = f"dialogues_{job_id}.json"
        else:
            download_name = f"results_{job_id}.json"

        return send_file(
            results_file,
            as_attachment=True,
            download_name=download_name,
            mimetype='application/json'
        )

    except Exception as e:
        logger.error(f"Error downloading job results: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/reuse-job/<job_id>', methods=['GET'])
def api_reuse_job_parameters(job_id):
    try:
        record = history_manager.get_record(job_id)
        if not record:
            return jsonify({"error": "Job not found"}), 404

        parameters = record.get('parameters', {})

        return jsonify({
            "status": "success",
            "job_id": job_id,
            "job_type": record.get('job_type'),
            "parameters": parameters
        })

    except Exception as e:
        logger.error(f"Error retrieving job parameters: {str(e)}")
        return jsonify({"error": str(e)}), 500


@app.route('/health', methods=['GET'])
def health_check():
    status = {
        "status": "healthy",
        "version": "1.0.0",
        "upload_folder": app.config['UPLOAD_FOLDER'],
        "results_folder": app.config['RESULTS_FOLDER'],
        "history_folder": app.config['HISTORY_FOLDER'],
        "upload_folder_writable": ensure_directory(app.config['UPLOAD_FOLDER']),
        "results_folder_writable": ensure_directory(app.config['RESULTS_FOLDER']),
        "history_folder_writable": ensure_directory(app.config['HISTORY_FOLDER']),
        "temp_dir_writable": ensure_directory(tempfile.gettempdir()),
        "optional_dependencies": OPTIONAL_DEPS_AVAILABLE
    }
    return jsonify(status)


def main():
    parser = argparse.ArgumentParser(description='Document Analysis and Dialogue Generation Tool')
    subparsers = parser.add_subparsers(dest='command', help='Command')

    parser_server = subparsers.add_parser('server', help='Start web server')
    parser_server.add_argument('--host', default='0.0.0.0', help='Server host address')
    parser_server.add_argument('--port', type=int, default=5233, help='Server port')
    parser_server.add_argument('--debug', action='store_true', help='Enable debug mode')

    parser_document = subparsers.add_parser('document', help='Process document into segments')
    parser_document.add_argument('file_path', help='Path to document file')
    parser_document.add_argument('--output', help='Output file path for results (JSON format)')
    parser_document.add_argument('--output-xlsx', help='Output file path for results (XLSX format)')
    parser_document.add_argument('--languages', nargs='+', default=['ch_sim', 'en'],
                                 help='OCR languages (e.g., ch_sim en ja)')
    parser_document.add_argument('--segment-size', type=int, default=1000,
                                 help='Maximum characters per segment')
    parser_document.add_argument('--no-overlap', action='store_true',
                                 help='Disable segment overlap (enabled by default)')
    parser_document.add_argument('--overlap-limit', type=int, default=200,
                                 help='Maximum character overlap')
    parser_document.add_argument('--no-clean-ai', action='store_true',
                                 help='Disable AI training data cleaning (enabled by default)')
    parser_document.add_argument('--replace-whitespace', action='store_true',
                                 help='Replace consecutive whitespace with single space')
    parser_document.add_argument('--remove-urls-emails', action='store_true',
                                 help='Remove all URLs and email addresses')
    parser_document.add_argument('--disable-ocr', action='store_true',
                                 help='Disable OCR for image processing')

    parser_dialogue = subparsers.add_parser('dialogue', help='Generate dialogues from text chunks')
    parser_dialogue.add_argument('input_file', help='JSON file containing text chunks')
    parser_dialogue.add_argument('--api-url', help='Base API URL')
    parser_dialogue.add_argument('--api-key', help='API authentication key')
    parser_dialogue.add_argument('--no-api', action='store_true', help='Generate placeholders without API')
    parser_dialogue.add_argument('--no-api-key', action='store_true', help='Use API without requiring an API key')
    parser_dialogue.add_argument('--output-dir', default='dialogues', help='Directory to save dialogue results')
    parser_dialogue.add_argument('--model', default='gpt-4o-mini', help='Model to use')
    parser_dialogue.add_argument('--rounds', type=int, default=3, help='Number of dialogue rounds')
    parser_dialogue.add_argument('--max-workers', type=int, default=4, help='Maximum parallel worker threads')
    parser_dialogue.add_argument('--language', default='en-zh', choices=['en', 'zh', 'en-zh'],
                                 help='Language mode: English (en), Chinese (zh), or bilingual (en-zh)')
    parser_dialogue.add_argument('--max-generation-mode', action='store_true',
                                 help='Enable maximum generation mode')
    parser_dialogue.add_argument('--max-questions', type=int, default=10,
                                 help='Maximum questions to generate')
    parser_dialogue.add_argument('--token-limit', type=int, default=1000,
                                 help='Maximum tokens per dialogue round')

    parser_integrated = subparsers.add_parser('integrated', help='Run integrated document and dialogue workflow')
    parser_integrated.add_argument('file_path', help='Path to document file')
    parser_integrated.add_argument('--api-url', help='Base API URL')
    parser_integrated.add_argument('--api-key', help='API authentication key')
    parser_integrated.add_argument('--no-api', action='store_true', help='Generate placeholders without API')
    parser_integrated.add_argument('--no-api-key', action='store_true', help='Use API without requiring an API key')
    parser_integrated.add_argument('--output-dir', default='results', help='Directory to save results')
    parser_integrated.add_argument('--model', default='gpt-4o-mini', help='Model to use')
    parser_integrated.add_argument('--dialogue-rounds', type=int, default=3, help='Number of dialogue rounds')
    parser_integrated.add_argument('--max-workers', type=int, default=4, help='Maximum parallel worker threads')
    parser_integrated.add_argument('--language', default='en-zh', choices=['en', 'zh', 'en-zh'],
                                   help='Language mode: English (en), Chinese (zh), or bilingual (en-zh)')
    parser_integrated.add_argument('--max-generation-mode', action='store_true',
                                   help='Enable maximum generation mode')
    parser_integrated.add_argument('--max-questions', type=int, default=5,
                                   help='Maximum questions to generate per segment')
    parser_integrated.add_argument('--token-limit', type=int, default=1000,
                                   help='Maximum tokens per dialogue round')
    parser_integrated.add_argument('--languages', nargs='+', default=['ch_sim', 'en'],
                                   help='OCR languages (e.g., ch_sim en ja)')
    parser_integrated.add_argument('--segment-size', type=int, default=1000,
                                   help='Maximum characters per segment')
    parser_integrated.add_argument('--no-overlap', action='store_true',
                                   help='Disable segment overlap (enabled by default)')
    parser_integrated.add_argument('--overlap-limit', type=int, default=200,
                                   help='Maximum character overlap')
    parser_integrated.add_argument('--no-clean-ai', action='store_true',
                                   help='Disable AI training data cleaning (enabled by default)')
    parser_integrated.add_argument('--disable-ocr', action='store_true',
                                   help='Disable OCR for image processing')

    args = parser.parse_args()

    if args.command == 'server':
        port = int(os.environ.get('PORT', args.port))
        host = os.environ.get('HOST', args.host)
        debug = os.environ.get('DEBUG', str(args.debug)).lower() == 'true'

        logger.info(f"Starting server on {host}:{port} (debug={debug})")
        app.run(host=host, port=port, debug=debug)

    elif args.command == 'document':
        try:
            segments = process_document(
                args.file_path,
                languages=args.languages,
                segment_size=args.segment_size,
                overlap=not args.no_overlap,
                overlap_limit=args.overlap_limit,
                clean_for_ai=not args.no_clean_ai,
                replace_whitespace=args.replace_whitespace,
                remove_urls_emails=args.remove_urls_emails,
                disable_ocr=args.disable_ocr
            )

            print(f"Document processed into {len(segments)} segments")

            if args.output:
                with open(args.output, 'w', encoding='utf-8') as f:
                    json.dump(segments, f, ensure_ascii=False, indent=2)
                print(f"Results saved to JSON: {args.output}")

            if args.output_xlsx and OPTIONAL_DEPS_AVAILABLE:
                export_segments_to_xlsx(segments, args.output_xlsx)
                print(f"Results saved to XLSX: {args.output_xlsx}")
            elif args.output_xlsx and not OPTIONAL_DEPS_AVAILABLE:
                print("Warning: Excel export requires openpyxl dependency. XLSX export skipped.")

            if not args.output and not args.output_xlsx:
                for i, segment in enumerate(segments[:3]):
                    print(f"Segment {i + 1}:")
                    print(f"Type: {segment.get('type', 'unknown')}")
                    content_preview = segment['content'][:100] + "..." if len(segment['content']) > 100 else segment[
                        'content']
                    print(f"Content: {content_preview}")
                    print("-" * 50)

                if len(segments) > 3:
                    print(f"...and {len(segments) - 3} more segments")

        except Exception as e:
            logger.error(f"Error processing document: {str(e)}")
            import traceback
            traceback.print_exc()

    elif args.command == 'dialogue':
        try:
            with open(args.input_file, 'r', encoding='utf-8') as f:
                data = json.load(f)

            if isinstance(data, list):
                chunks = [item.get('content', '') for item in data]
            else:
                if 'segments' in data:
                    chunks = [item.get('content', '') for item in data['segments']]
                else:
                    chunks = list(data.values())

            print(f"Loaded {len(chunks)} text chunks for dialogue generation")

            use_api = not args.no_api
            require_api_key = not args.no_api_key

            if use_api and not args.api_url:
                print("Error: API URL is required when using API mode.")
                print("Use --no-api flag to run without API.")
                return

            if use_api and require_api_key and not args.api_key:
                print("Error: API key is required when using API with key requirement.")
                print("Use --no-api-key flag to use API without requiring API key.")
                return

            generator = ParallelDialogueGenerator(
                api_url=args.api_url,
                api_key=args.api_key,
                output_dir=args.output_dir,
                model=args.model,
                max_generation_mode=args.max_generation_mode,
                max_questions=args.max_questions,
                dialogue_token_limit=args.token_limit,
                language=args.language,
                use_api=use_api,
                require_api_key=require_api_key
            )

            results = generator.process_chunks_parallel(
                chunks=chunks,
                rounds=args.rounds,
                max_workers=args.max_workers
            )

            print(f"Generated dialogues for {len(results)} chunks")
            overall_results_path = os.path.join(args.output_dir, 'all_dialogues.json')
            with open(overall_results_path, 'w', encoding='utf-8') as f:
                json.dump(results, f, ensure_ascii=False, indent=2)

            print(f"All dialogues saved to {overall_results_path}")

        except Exception as e:
            logger.error(f"Error generating dialogues: {str(e)}")
            import traceback
            traceback.print_exc()

    elif args.command == 'integrated':
        try:
            use_api = not args.no_api
            require_api_key = not args.no_api_key

            if use_api and not args.api_url:
                print("Error: API URL is required when using API mode.")
                print("Use --no-api flag to run without API.")
                return

            if use_api and require_api_key and not args.api_key:
                print("Error: API key is required when using API with key requirement.")
                print("Use --no-api-key flag to use API without requiring API key.")
                return

            result = process_document_and_generate_dialogues(
                file_path=args.file_path,
                api_url=args.api_url,
                api_key=args.api_key,
                model=args.model,
                dialogue_rounds=args.dialogue_rounds,
                max_workers=args.max_workers,
                output_dir=args.output_dir,
                language=args.language,
                max_generation_mode=args.max_generation_mode,
                max_questions=args.max_questions,
                dialogue_token_limit=args.token_limit,
                segment_size=args.segment_size,
                overlap=not args.no_overlap,
                overlap_limit=args.overlap_limit,
                clean_for_ai=not args.no_clean_ai,
                languages=args.languages,
                use_api=use_api,
                disable_ocr=args.disable_ocr
            )

            if result["status"] == "success":
                print(f"Integration workflow completed successfully:")
                print(f"- Document processed into {result['segments_count']} segments")
                print(f"- Generated dialogues for {result['dialogues_count']} segments")
                print(f"- Results saved to {result['results_file']}")
                if result.get('xlsx_file'):
                    print(f"- Excel format saved to {result['xlsx_file']}")
            else:
                print(f"Integration workflow failed: {result['error']}")

        except Exception as e:
            logger.error(f"Error in integrated workflow: {str(e)}")
            import traceback
            traceback.print_exc()

    else:
        parser.print_help()

if __name__ == "__main__":
    main()